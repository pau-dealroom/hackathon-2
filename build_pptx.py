"""Generate a one-pager PowerPoint deck for the hackathon."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Palette (matches Dealroom-themed HTML one-pager) ─────────────────────────
NAVY        = RGBColor(0x0F, 0x1A, 0x2E)
TEXT        = RGBColor(0x0F, 0x1A, 0x2E)
MUTED       = RGBColor(0x6B, 0x72, 0x80)
MUTED2      = RGBColor(0x9C, 0xA3, 0xAF)
BORDER      = RGBColor(0xE5, 0xE7, 0xEB)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT      = RGBColor(0x36, 0x62, 0xE3)
ACCENT_LITE = RGBColor(0xEE, 0xF2, 0xFF)
GREEN       = RGBColor(0x16, 0xA3, 0x4A)
GREEN_LITE  = RGBColor(0xF0, 0xFD, 0xF4)
BG          = RGBColor(0xF5, 0xF6, 0xFA)
LIGHT_GREY  = RGBColor(0xF3, 0xF4, 0xF6)


def add_text(shape, text, size, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
             font="Inter", line_spacing=1.2):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tf


def fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def no_line(shape):
    shape.line.fill.background()


def border(shape, rgb, weight=0.75):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(weight)


# ── Build slide ──────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
fill(bg, BG)
no_line(bg)

# White card
card_left, card_top = Inches(0.6), Inches(0.5)
card_w, card_h      = Inches(12.13), Inches(6.5)
card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_left, card_top, card_w, card_h)
card.adjustments[0] = 0.04
fill(card, WHITE)
border(card, BORDER, 1)

# ── Top bar ──────────────────────────────────────────────────────────────────
brand_y = Inches(0.95)
mark = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), brand_y, Inches(0.42), Inches(0.42))
mark.adjustments[0] = 0.25
fill(mark, NAVY)
no_line(mark)

# Bar chart in mark — three small white bars
bar_specs = [(0.10, 0.27, 0.07, 0.13), (0.18, 0.20, 0.07, 0.20), (0.26, 0.13, 0.07, 0.27)]
for bx, by, bw, bh in bar_specs:
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(1.0 + bx), Inches(brand_y.inches + by),
        Inches(bw), Inches(bh))
    fill(b, WHITE)
    no_line(b)

brand_text = slide.shapes.add_textbox(Inches(1.52), Inches(0.98), Inches(2.5), Inches(0.36))
add_text(brand_text, "dealroom.co", size=15, bold=True, color=TEXT)

# Badge (top right)
badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(10.95), Inches(0.98), Inches(1.7), Inches(0.36))
badge.adjustments[0] = 0.5
fill(badge, ACCENT_LITE)
border(badge, ACCENT, 0.6)
add_text(badge, "Hackathon · 2026", size=10, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER)
badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# Top divider
div1 = slide.shapes.add_connector(1, Inches(1.0), Inches(1.55), Inches(12.7), Inches(1.55))
div1.line.color.rgb = BORDER
div1.line.width = Pt(0.5)

# ── Title block ──────────────────────────────────────────────────────────────
eyebrow = slide.shapes.add_textbox(Inches(1.0), Inches(1.85), Inches(6), Inches(0.3))
add_text(eyebrow, "PROJECT · ONE-PAGER", size=10, bold=True, color=ACCENT)

title = slide.shapes.add_textbox(Inches(1.0), Inches(2.15), Inches(11.5), Inches(0.7))
add_text(title, "Custom page for VCs", size=36, bold=True, color=TEXT)

goal = slide.shapes.add_textbox(Inches(1.0), Inches(2.95), Inches(11), Inches(0.95))
tf = goal.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_right = Emu(0)
p = tf.paragraphs[0]
p.line_spacing = 1.4
r1 = p.add_run(); r1.text = "Goal: "; r1.font.bold = True
r1.font.name = "Inter"; r1.font.size = Pt(14); r1.font.color.rgb = TEXT
r2 = p.add_run()
r2.text = ("spark curiosity through personalized content with the goal of wanting to see more "
           "— only sending companies not already in their portfolio.")
r2.font.name = "Inter"; r2.font.size = Pt(14); r2.font.color.rgb = MUTED

# ── Criteria label ───────────────────────────────────────────────────────────
crit_label = slide.shapes.add_textbox(Inches(1.0), Inches(4.1), Inches(6), Inches(0.3))
add_text(crit_label, "EVALUATION CRITERIA", size=10, bold=True, color=MUTED2)

# ── 3 cards ──────────────────────────────────────────────────────────────────
card_specs = [
    ("Usefulness",  "✓ Yes",   "Saves the sales team manual research time and sparks curiosity in the VC prospect — every page is a unique conversation starter."),
    ("Working demo","✓ Yes",   "Live Flask app generating bespoke Company & LP recommendation pages end-to-end from a single VC name or URL input."),
    ("Learning leap","✓ Loads","Learned a lot — working with API connectors, fixing bugs across the stack, prompt design for GPT-driven thesis extraction, and navigating the next-gen Dealroom API."),
]

cards_top = Inches(4.5)
card_h    = Inches(1.95)
card_w    = Inches(3.75)
gap       = Inches(0.18)
start_x   = Inches(1.0)

for i, (ctitle, pill_text, body) in enumerate(card_specs):
    x = start_x + (card_w + gap) * i
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, cards_top, card_w, card_h)
    c.adjustments[0] = 0.05
    fill(c, WHITE)
    border(c, BORDER, 0.75)

    # Left accent stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        x, cards_top, Inches(0.05), card_h)
    fill(stripe, ACCENT)
    no_line(stripe)

    # Title
    t = slide.shapes.add_textbox(x + Inches(0.25), cards_top + Inches(0.18), Inches(2.4), Inches(0.32))
    add_text(t, ctitle, size=13, bold=True, color=TEXT)

    # Pill
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        x + Inches(2.7), cards_top + Inches(0.2), Inches(0.85), Inches(0.27))
    pill.adjustments[0] = 0.5
    fill(pill, GREEN_LITE)
    no_line(pill)
    add_text(pill, pill_text, size=9, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    pill.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Body
    body_box = slide.shapes.add_textbox(x + Inches(0.25), cards_top + Inches(0.62),
                                        card_w - Inches(0.5), card_h - Inches(0.75))
    bf = body_box.text_frame
    bf.word_wrap = True
    bf.margin_left = bf.margin_right = Emu(0)
    bp = bf.paragraphs[0]
    bp.line_spacing = 1.4
    br = bp.add_run()
    br.text = body
    br.font.name = "Inter"; br.font.size = Pt(11); br.font.color.rgb = MUTED

# ── Footer ───────────────────────────────────────────────────────────────────
foot_top = Inches(6.75)
div2 = slide.shapes.add_connector(1, Inches(1.0), foot_top, Inches(12.7), foot_top)
div2.line.color.rgb = BORDER
div2.line.width = Pt(0.5)

tag_specs = ["Dealroom API", "Flask", "GPT-4o", "Python"]
tag_x = Inches(1.0)
for tname in tag_specs:
    tw = Inches(0.85 + 0.05 * len(tname))
    tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        tag_x, foot_top + Inches(0.12), tw, Inches(0.28))
    tag.adjustments[0] = 0.35
    fill(tag, LIGHT_GREY)
    border(tag, BORDER, 0.5)
    add_text(tag, tname, size=9, color=MUTED, align=PP_ALIGN.CENTER)
    tag.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    tag_x = tag_x + tw + Inches(0.08)

repo = slide.shapes.add_textbox(Inches(9.5), foot_top + Inches(0.15), Inches(3.2), Inches(0.3))
add_text(repo, "github.com/pau-dealroom/hackathon-2", size=10,
         color=MUTED2, align=PP_ALIGN.RIGHT)

out = "hackathon_slide.pptx"
prs.save(out)
print(f"✓ {out}")
